import os
import streamlit as st
import tempfile
import pandas as pd
import json
import xml.etree.ElementTree as ET
import yaml
from bs4 import BeautifulSoup
from pptx import Presentation
from docx import Document

from langchain.document_loaders import PyPDFLoader, TextLoader
from langchain.indexes import VectorstoreIndexCreator
from langchain.chains import RetrievalQA
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate

from ibm_watson_machine_learning.foundation_models import Model
from ibm_watson_machine_learning.foundation_models.extensions.langchain import WatsonxLLM
from ibm_watson_machine_learning.metanames import GenTextParamsMetaNames as GenParams
from ibm_watson_machine_learning.foundation_models.utils.enums import DecodingMethods

# Initialize index and chain to None
index = None
rag_chain = None

# Custom loader for DOCX files
class DocxLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        document = Document(self.file_path)
        text_content = [para.text for para in document.paragraphs]
        return " ".join(text_content)

# Custom loader for PPTX files
class PptxLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        presentation = Presentation(self.file_path)
        text_content = [shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")]
        return " ".join(text_content)

# Custom loader for additional file types
def load_csv(file_path):
    df = pd.read_csv(file_path)
    page_size = 100
    page_number = st.number_input("Page number", min_value=1, max_value=(len(df) // page_size) + 1, step=1, value=1)
    start_index = (page_number - 1) * page_size
    end_index = start_index + page_size
    st.dataframe(df.iloc[start_index:end_index])
    return df.to_string(index=False)

def load_json(file_path):
    with open(file_path, 'r') as file:
        data = json.load(file)
    return json.dumps(data, indent=2)

def load_xml(file_path):
    tree = ET.parse(file_path)
    root = tree.getroot()
    return ET.tostring(root, encoding="unicode")

def load_yaml(file_path):
    with open(file_path, 'r') as file:
        data = yaml.safe_load(file)
    return yaml.dump(data)

def load_html(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        soup = BeautifulSoup(file, 'html.parser')
    return soup.get_text()

# Caching function to load various file types
@st.cache_resource
def load_file(file_name, file_type):
    loaders = []
    text = None

    if file_type == "pdf":
        loaders = [PyPDFLoader(file_name)]
    elif file_type == "docx":
        loader = DocxLoader(file_name)
        text = loader.load()
    elif file_type == "pptx":
        loader = PptxLoader(file_name)
        text = loader.load()
    elif file_type == "txt":
        loaders = [TextLoader(file_name)]
    elif file_type == "csv":
        text = load_csv(file_name)
    elif file_type == "json":
        text = load_json(file_name)
    elif file_type == "xml":
        text = load_xml(file_name)
    elif file_type == "yaml":
        text = load_yaml(file_name)
    elif file_type == "html":
        text = load_html(file_name)
    elif file_type == "htm":
        text = load_html(file_name)    
    else:
        st.error("Unsupported file type.")
        return None

    if text:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".txt") as temp_file:
            temp_file.write(text.encode("utf-8"))
            temp_file_path = temp_file.name
        loaders = [TextLoader(temp_file_path)]

    if loaders:
        index = VectorstoreIndexCreator(
            embedding=HuggingFaceEmbeddings(model_name="all-MiniLM-L12-v2"),
            text_splitter=RecursiveCharacterTextSplitter(chunk_size=450, chunk_overlap=50)
        ).from_loaders(loaders)
        st.success("Index created successfully!")
        return index
    return None

# Watsonx API setup
watsonx_api_key =  os.getenv("WATSONX_API_KEY")
watsonx_project_id = os.getenv("WATSONX_PROJECT_ID")

if not watsonx_api_key or not watsonx_project_id:
    st.error("API Key or Project ID is not set. Please set them as environment variables.")

prompt_template_br = PromptTemplate(
    input_variables=["context", "question"],
    template="""<|begin_of_text|><|start_header_id|>system<|end_header_id|>
I am a helpful assistant.

<|eot_id|>
{context}
<|start_header_id|>user<|end_header_id|>
{question}<|eot_id|>
"""
)

# Sidebar tab selection
with st.sidebar:
    st.title("Watsonx RAG: Multi-Document Retrieval")
    tab = st.radio("Select Mode", ["Watsonx Discovery", "File Upload"])

# Handle Watsonx Discovery tab
if tab == "Watsonx Discovery":
    # Watsonx Discovery setup
    st.info("Ask a question using Watsonx Discovery")
    uploaded_file = st.text_input("Enter document ID from Watsonx Discovery")
    
    if uploaded_file:
        # Use Watsonx Discovery to fetch documents and process
        # Replace with your Discovery query and model call
        
        st.info(f"Querying Watsonx Discovery with document ID: {uploaded_file}")
        # Assuming you have a method to query Discovery based on the doc ID
        # Replace `fetch_from_discovery` with actual API interaction

        # Assuming the response comes as a list of documents or text
        context = "Sample response from Watsonx Discovery based on the doc ID."
        question = st.text_input("Ask your question about the document")

        if question:
            st.info(f"Question: {question}")
            # LLM call to Watsonx for question-answering
            response = "Answer from LLM based on Discovery content."
            st.write(response)
        
# Handle File Upload tab
elif tab == "File Upload":
    with st.sidebar:
        watsonx_model = st.selectbox("Model", ["meta-llama/llama-3-405b-instruct", "codellama/codellama-34b-instruct-hf", "ibm/granite-20b-multilingual"])
        max_new_tokens = st.slider("Max output tokens", min_value=100, max_value=4000, value=600, step=100)
        decoding_method = st.radio("Decoding", (DecodingMethods.GREEDY.value, DecodingMethods.SAMPLE.value))

    # File upload logic
    uploaded_file = st.file_uploader("Upload a file", accept_multiple_files=False, type=["pdf", "docx", "txt", "pptx", "csv", "json", "xml", "yaml", "html"])

    if uploaded_file is not None:
        bytes_data = uploaded_file.read()
        st.write("Filename:", uploaded_file.name)
        with open(uploaded_file.name, 'wb') as f:
            f.write(bytes_data)
        file_type = uploaded_file.name.split('.')[-1].lower()
        index = load_file(uploaded_file.name, file_type)

    model_name = watsonx_model

    # Watsonx Model Setup (similar to your existing code)

    # Other code related to file handling and LLM interaction remains the same...
